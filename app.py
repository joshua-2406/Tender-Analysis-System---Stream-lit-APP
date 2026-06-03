import os
os.environ["PP_OCR_LOG_LEVEL"] = "ERROR"

import streamlit as st

# Page config
st.set_page_config(page_title="Professional Tender Presentation Generator", layout="wide")

import pdfplumber
import docx
import re
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import requests
from geopy.geocoders import Nominatim
from pptx.enum.text import MSO_AUTO_SIZE
import folium
from streamlit_folium import st_folium
from geopy.distance import geodesic
import pandas as pd
import io
import requests
from io import BytesIO
from PIL import ImageEnhance, ImageFilter
from pdf2image import convert_from_bytes
import pytesseract
from paddleocr import PaddleOCR

def add_footer(slide):
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = footer_box.text_frame
    tf.text = f"Generated on {datetime.now().strftime('%d %B %Y')} | TETAKISU Confidential"
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
# PaddleOCR import & init (after st.set_page_config!)
try:
    from paddleocr import PaddleOCR
    PADDLE_OCR_AVAILABLE = True
except ImportError:
    PADDLE_OCR_AVAILABLE = False

ocr_engine = None
if PADDLE_OCR_AVAILABLE:
    try:
        # ONLY use the NEW parameter — NO use_angle_cls
        ocr_engine = PaddleOCR(
            use_textline_orientation=True,   # Correct modern way
            lang='en',
            #show_log=False                   # Less spam
        )
    except Exception as e:
        ocr_engine = None
        print(f"PaddleOCR init error: {e}")  # Print to console only

# Now safe to show UI
st.title("Professional Tender Presentation Generator")

if PADDLE_OCR_AVAILABLE and ocr_engine:
    st.success("PaddleOCR initialized successfully! Using advanced OCR for scanned PDFs.")
else:
    st.warning("PaddleOCR not available. Install with: pip install paddleocr paddlepaddle\nFalling back to basic OCR.")

# ---- SAFE INITIALIZATION (FIXES NameError) ----
content_by_topic = {}
site_location = "N/A"
coords = None
climate = {}
infra = {}

# Upload multiple files
uploaded_files = st.file_uploader(
    "Upload Tender Files (PDF/DOCX/XLSX/XLS)",
    type=['pdf', 'docx', 'xlsx', 'xls'],
    accept_multiple_files=True,
    help="Upload all tender documents. Multiple files supported."
)

# Standard slide titles
STANDARD_SLIDES = [
    "Scope of Work",
    "Pre-Qualification Criteria & Technical Capability",
    "Technical Specifications",
    "Commercial Terms",
    "Documents Required",
    "Eligibility Criteria",
    "Payment Terms",
    "Security Deposit & Performance Bank Guarantee",
    "Liquidated Damages",
    "Warranty & Guarantee",
    "Safety Requirements",
    "Site Location & Details",
    "Climatic Conditions Summary",
    "Nearby Infrastructure",
    "Bid Schedule & Key Dates",
    "Delivery Period",
    "Quality Standards",
    "Evaluation Criteria",
    "Expected Competition & Market Analysis",
    "Bid Strategy & Recommendation",
    "Critical Provisions & Compliance",
    "Terms & Conditions Acceptance",
    "Make in India & Local Content",
    "Risk Analysis",
    "Risk Mitigation Measures",
    "Financial Analysis",
    "Organization & Team Structure",
    "Quality Assurance Plan",
    "Manufacturing Process Overview",
    "Inspection & Testing Protocol",
    "Packing & Transportation",
    "Key Success Factors",
    "Conclusion & Recommendation"
]

def summarize_text(raw_text, max_sentences=5):
    """
    Dynamically creates a short, simple summary from raw extracted text.
    - Removes duplicates/garbled parts
    - Takes top sentences
    - Makes it readable
    """
    # Clean text
    clean = re.sub(r'(\b\w+\b)(\s*\1){2,}', r'\1', raw_text)  # Remove repeats
    clean = re.sub(r'\s+', ' ', clean).strip()

    # Split into sentences
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', clean)

    # Simple extractive summary: Take first few unique sentences
    unique_sentences = []
    for s in sentences:
        s = s.strip()
        if s and len(s) > 20 and s not in unique_sentences:
            unique_sentences.append(s)
        if len(unique_sentences) >= max_sentences:
            break

    # Make bullet points
    summary = "Short Summary:\n" + "\n".join([f"- {s}" for s in unique_sentences if s])

    if not unique_sentences:
        summary = "No clear details found in this section."

    return summary

def detect_location_and_coords(text, tender_info):
    geolocator = Nominatim(user_agent="tender_locator_v5")

    # Step 1: Stronger regex extraction with more patterns (handles variations in tenders)
    patterns = [
        r"(?i)(?:site location|project site|work site|plant location|execution at|delivery at|located at|situated at|place of work)[\s:–\-]*([A-Za-z0-9\s\-/(),]+?)(?:,|\.|Tamil Nadu|District|India|$)",
        r"(?i)(?:kalpakkam|igcar|pfbr|ranipet|bhel|boiler auxiliaries|mahabalipuram|sadras|pudupattinam)[A-Za-z\s,]*?(?:tamil nadu)?",  # Optional: Boost for common ones, but not hardcoded
    ]

    candidates = []
    for pat in patterns:
        matches = re.findall(pat, text)
        for m in matches:
            if isinstance(m, tuple): m = m[0]  # Handle grouped matches
            loc = re.sub(r'\s+', ' ', m).strip()
            if len(loc) > 5 and loc.lower() not in ['india', 'tamil nadu', 'bharat', 'heavy', 'electricals', 'limited']:
                candidates.append(loc)

    # Add from tender_info if available
    if tender_info.get("location") and tender_info["location"] != "N/A":
        candidates.insert(0, tender_info["location"])

    candidates = list(dict.fromkeys(candidates))  # Remove duplicates

    # Step 2: For each candidate, build contextual queries and auto-fetch from Nominatim (dynamic source)
    for place in candidates:
        context = tender_info.get('client', '') or tender_info.get('tender_name', '') or ''
        queries = [
            f"{place}, India",
            f"{place}, {context}, India",
            f"{place}, Tamil Nadu, India" if 'tamil nadu' in text.lower() else f"{place}, India",
            f"{place} project site, India",
            f"{place} industrial area, India",
        ]

        for q in queries:
            try:
                geo = geolocator.geocode(q, timeout=10)
                if geo and 8 < geo.latitude < 38 and 68 < geo.longitude < 98:  # Strict India geographic bounds to avoid wrong places
                    # Return cleaned name (from source) + coords
                    clean_name = geo.address.split(',')[0].strip().title()
                    return clean_name, (geo.latitude, geo.longitude)
            except:
                continue  # Skip failed queries silently

    # Step 3: Dynamic fallback - extract state/city keywords and query again
    state_match = re.search(
        r"(?i)(Tamil Nadu|Karnataka|Maharashtra|Odisha|Gujarat|Rajasthan|Uttar Pradesh|Telangana|Andhra Pradesh|Kerala|West Bengal|Bihar|Chhattisgarh|Delhi|Mumbai|Hyderabad|Bangalore|Chennai)",
        text
    )
    if state_match:
        fallback_place = state_match.group(1).title()
        try:
            geo = geolocator.geocode(f"{fallback_place} industrial area, India", timeout=10)
            if geo:
                clean_name = geo.address.split(',')[0].strip().title()
                return clean_name, (geo.latitude, geo.longitude)
        except:
            pass

    # Absolute last resort (rarely reached)
    return "Location Not Identified (check tender text)", None


@st.cache_data(show_spinner=False)
def get_climate_data(lat, lon):
    try:
        url = (
            "https://api.open-meteo.com/v1/forecast"
            f"?latitude={lat}&longitude={lon}"
            f"&current=temperature_2m,relative_humidity_2m,precipitation"
            f"&timezone=Asia/Kolkata"  # Ensures correct local time
            f"&forecast_days=1"
        )
        response = requests.get(url, timeout=10)
        if response.status_code != 200:
            return {"temperature": "N/A", "humidity": "N/A", "precipitation": "N/A"}
        
        data = response.json()
        current = data.get("current", {})
        
        temp = current.get("temperature_2m")
        hum = current.get("relative_humidity_2m")
        rain = current.get("precipitation")
        
        return {
            "temperature": f"{temp} °C" if temp is not None else "N/A",
            "humidity": f"{hum} %" if hum is not None else "N/A",
            "precipitation": f"{rain} mm" if rain is not None else "N/A"
        }
    except Exception as e:
        return {"temperature": "N/A", "humidity": "N/A", "precipitation": "N/A"}
    

@st.cache_data(show_spinner=False)
def get_nearby_infra(lat, lon):
    infra = {
        "Railway Station": [],
        "Airport": [],
        "Bus Stand": [],
        "Hotel": []
    }

    tags = {
        "Railway Station": [
            'railway=station',
            'railway=halt'
        ],
        "Airport": [
            'aeroway=airport',
            'aeroway=aerodrome'
        ],
        "Bus Stand": [
            'amenity=bus_station',
            'highway=bus_stop'
        ],
        "Hotel": [
            'tourism=hotel',
            'tourism=guest_house'
        ],
    }

    radii = {
    "Railway Station": 40000,   # 40 km
    "Airport": 150000,          # 150 km 
    "Bus Stand": 30000,
    "Hotel": 30000
}

    for infra_type, tag_list in tags.items():
        collected = []

        for tag in tag_list:
            query = f"""
            [out:json][timeout:25];
            (
              node[{tag}](around:{radii[infra_type]},{lat},{lon});
              way[{tag}](around:{radii[infra_type]},{lat},{lon});
              relation[{tag}](around:{radii[infra_type]},{lat},{lon});
            );
            out center;
            """

            try:
                resp = requests.get(
                    "https://overpass-api.de/api/interpreter",
                    params={'data': query},
                    timeout=30
                ).json()

                for el in resp.get("elements", []):
                    name = el.get("tags", {}).get("name")
                    if not name:
                        continue

                    clat = el.get("lat") or el.get("center", {}).get("lat")
                    clon = el.get("lon") or el.get("center", {}).get("lon")

                    if clat and clon:
                        dist = geodesic((lat, lon), (clat, clon)).km
                        collected.append((name, dist))
            except:
                pass

        # Sort & keep nearest 3 unique
        unique = {}
        for name, dist in sorted(collected, key=lambda x: x[1]):
            if name not in unique:
                unique[name] = round(dist, 1)
            if len(unique) == 3:
                break

        infra[infra_type] = [
            f"{name} ({dist} km)" for name, dist in unique.items()
        ]

    return infra


def extract_text_and_tables_from_pdf(file_bytes, filename):
    text = ""
    tables = []

    try:
        # Step 1: Try normal PDF extraction (fast for searchable PDFs)
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
                page_tables = page.extract_tables()
                for t in page_tables:
                    if t and any(any(cell for cell in row if cell) for row in t):
                        df = pd.DataFrame(t[1:], columns=t[0]) if t else pd.DataFrame()
                        tables.append(df)

        # Step 2: Check if OCR is needed (low text or known garbage)
        if len(text.strip()) < 150 or "nokh" in text.lower() or "nan" in text.lower() or len(text.split()) < 50:
            st.info(f"{filename} appears scanned or garbled → using advanced OCR...")
            text = ""  # Reset text - we will rebuild with OCR
            images = convert_from_bytes(file_bytes, dpi=350)

            # Try PaddleOCR first (advanced)
            if PADDLE_OCR_AVAILABLE and ocr_engine:
                for img in images:
                    try:
                        result = ocr_engine.ocr(img)
                        page_text = ""
                        if result and len(result) > 0 and result[0]:
                            for line in result[0]:
                                if line and len(line) >= 2 and line[1]:
                                    page_text += line[1][0] + "\n"
                        text += page_text + "\n"
                    except Exception as ocr_err:
                        text += f"[PaddleOCR Error on page: {str(ocr_err)}]\n"

            # Fallback to enhanced pytesseract if PaddleOCR fails/unavailable
            else:
                st.warning("PaddleOCR not available → using enhanced pytesseract fallback")
                for img in images:
                    try:
                        img = img.convert('L')
                        enhancer = ImageEnhance.Contrast(img)
                        img = enhancer.enhance(2.5)
                        img = img.filter(ImageFilter.MedianFilter(size=3))

                        ocr_text = pytesseract.image_to_string(
                            img,
                            lang='eng',
                            config='--oem 3 --psm 6 -c tessedit_char_whitelist=ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789.,()/-:;"'
                        )
                        text += ocr_text + "\n"
                    except Exception as tess_err:
                        text += f"[Tesseract Error on page: {str(tess_err)}]\n"

    except Exception as e:
        st.error(f"Error processing PDF {filename}: {e}")
        text = "[Processing Failed - Check file format]"

    # Step 3: Final cleanup of common OCR garbage
    text = re.sub(r'(\b\w+\b)(\s+\1){3,}', r'\1', text)           # Remove 3+ word repeats
    text = re.sub(r'(Nokh|Nalt|NAN|NaN|NOKH|NALT)\s*', '', text, flags=re.I)  # Remove junk words
    text = re.sub(r'\n{3,}', '\n\n', text)                        # Reduce multiple blank lines

    return text, tables

def extract_text_and_tables_from_docx(file_bytes, filename):
    text = ""
    tables = []
    try:
        doc = docx.Document(BytesIO(file_bytes))
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        for table in doc.tables:
            data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                data.append(row_data)
            if data:
                # Assume first row is header
                df = pd.DataFrame(data[1:], columns=data[0] if data[0] else None)
                tables.append(df)
    except Exception as e:
        st.error(f"Error processing DOCX {filename}: {e}")
        text = "[Unreadable DOCX]"
    
    return text, tables


def extract_text_and_tables_from_excel(file_bytes, filename):
    text = ""
    tables = []
    try:
        xl = pd.ExcelFile(BytesIO(file_bytes))
        for sheet in xl.sheet_names:
            df = pd.read_excel(xl, sheet_name=sheet, engine='openpyxl' if filename.endswith('.xlsx') else 'xlrd')
            text += f"Sheet: {sheet}\n{df.to_string(index=False)}\n\n"
            tables.append(df)
    except Exception as e:
        st.error(f"Error processing Excel {filename}: {e}")
        text = "[Unreadable Excel]"
    
    return text, tables

import re

def generate_table_summary(all_tables):
    if not all_tables:
        return "No tables detected in the uploaded documents."

    combined_lower = ""
    combined_original = ""
    for df in all_tables:
        try:
            table_str = df.to_string(index=False, header=True)
            combined_lower += table_str.lower() + "\n\n"
            combined_original += table_str + "\n\n"
        except:
            combined_lower += "table (format error)\n\n"
            combined_original += "table (format error)\n\n"

    summary_lines = ["Key insights from tables across all documents:"]

    # Broad categories (dynamic)
    categories = {
        "Sampling and inspection plans": ["sampling", "aql", "lot size", "sample size", "accept", "reject", "double sampling"],
        "Dimensional tolerances and deviations": ["tolerance", "deviation", "dimension", "linear", "angular", "±", "allowable", "iso"],
        "Surface preparation and painting": ["painting", "surface preparation", "dft", "primer", "finish coat", "blast", "sa ", "profile", "zinc", "enamel"],
        "Welding procedures and qualifications": ["welding", "wps", "pqr", "qualification", "welder", "electrode", "thickness range", "position"],
        "Manufacturing quality plans (MQP) and controls": ["mqp", "quality plan", "inprocess", "class a", "class b", "tc", "witness", "verification"],
        "Material specifications and testing": ["chemical composition", "heat treatment", "mechanical test", "ultrasonic", "ndt", "grain size", "killed", "normalized"],
        "Revision history and approvals": ["rev", "revision", "prepared by", "reviewed by", "approved by", "signature", "date"]
    }

    detected = [theme for theme, keywords in categories.items() if any(kw in combined_lower for kw in keywords)]

    if detected:
        summary_lines.append("• " + "\n• ".join(detected))
    else:
        summary_lines.append("• Tables contain general procedural and approval records")

    extras = []

    # DFT
    dft_matches = re.findall(r"dft[^\d]*(\d+\s*[\–\-–]\s*\d+\s*µm|\d+\s*µm)", combined_original, re.I)
    if dft_matches:
        extras.append(f"  - Dry Film Thickness (DFT): {', '.join(sorted(set(dft_matches)))}")

    # Surface profile
    profile_matches = re.findall(r"profile\s*([\d–\-]+\s*µm|35\s*[\–\-–]\s*50\s*µm)", combined_original, re.I)
    if profile_matches:
        extras.append(f"  - Surface profile: {', '.join(sorted(set(profile_matches)))}")

    # Lot sizes from sampling tables (improved regex to catch more formats like "Lot Size 2-1200", "lot 2 to 1200", "above 1201")
    lot_matches = re.findall(r"(?:lot size|lot|size)\s*[^\d]*([\d\s–\-]+(?:\s*(?:to|–|and|above)\s*[\d\s–\-]+)?)", combined_original, re.I)
    if lot_matches:
        unique_lots = sorted(set(l.strip() for l in lot_matches if l.strip()))
        extras.append(f"  - Lot sizes: {', '.join(unique_lots)}")

    # Inspection classes (A/B/C from MQP tables - improved to catch "Class A", "class B", etc.)
    class_matches = re.findall(r"class\s*([A-Ca-c])\b", combined_original, re.I)
    if class_matches:
        unique_classes = sorted(set(c.upper() for c in class_matches))
        extras.append(f"  - Inspection classes: Class {', '.join(unique_classes)} (A=Critical, B=Major, C=Minor)")

    # AQL
    aql_matches = re.findall(r"aql[^\d]*([\d\.]+)", combined_original, re.I)
    if aql_matches:
        extras.append(f"  - AQL levels: {', '.join(sorted(set(aql_matches)))}%")

    # Tolerances
    tol_matches = re.findall(r"±\s*([\d\.]+)\s*(mm|°)", combined_original)
    if tol_matches:
        extras.append(f"  - Tolerance ranges include ±{', '.join([m[0] for m in tol_matches[:5]])} {tol_matches[0][1] if tol_matches else ''}")

    # Revision count (more accurate: only Rev/Revision headers)
    rev_count = len(re.findall(r"(?:Rev|REVISION)\s*[0-9]+", combined_original, re.I))
    if rev_count > 0:
        extras.append(f"  - Multiple revisions tracked (up to Rev {rev_count})")

    # Class A/B/C
    class_matches = re.findall(r"class\s*([a-cA-C])", combined_original, re.I)
    if class_matches:
        extras.append(f"  - Inspection classes: Class {', '.join(sorted(set(c.upper() for c in class_matches)))} (A=Critical, B=Major, etc.)")

    if extras:
        summary_lines.append("\nNotable specifications extracted:")
        summary_lines.extend(extras)

    # Safety fallback
    if len(detected) <= 1 and "revision" in combined_lower:
        return "Tables primarily contain document revision history and approval records."

    return "\n".join(summary_lines) if summary_lines else "Tables contain basic procedural records."

    # === END OF NEW BLOCKS ===
    

def extract_tender_info(text):
    info = {
        'tender_no': 'N/A',
        'tender_name': 'Tender Analysis',
        'client': 'N/A',
        'location': 'N/A',
        'due_date': 'N/A',
        'nature_of_work': 'N/A'
    }

    # Tender Number
    m = re.search(r"(?i)(?:tender|enquiry|e-tender)\s*(?:no|id|number)\.?\s*[:\-]?\s*([\w\/\-]+)", text)
    if m: 
        info['tender_no'] = m.group(1).strip()
    
    # Tender Name / Nature of Work
    m = re.search(r"(?i)(?:tender enquiry for|contract for|e-tender.*?is a tender enquiry for)\s*(.*?components)", text, re.DOTALL)
    if m: 
        name = re.sub(r'\s+', ' ', m.group(1).strip()).title()
        info['tender_name'] = name
        info['nature_of_work'] = name
    
    if info['tender_name'] == 'Tender Analysis':
        m = re.search(r"(?i)(complete manufacturing|supply of|fabrication).*?(esp|components|structures)", text, re.I)
        if m: 
            info['tender_name'] = m.group(0).strip().title()
            info['nature_of_work'] = info['tender_name']
    
    # Client
    m = re.search(r"(?i)^([A-Z][A-Za-z\s]+Limited)", text, re.M)
    if m: 
        info['client'] = m.group(1).strip()
    
    if info['client'] == 'N/A':
        m = re.search(r"(?i)(issued by|client|authority):\s*([A-Za-z\s]+)", text)
        if m: 
            info['client'] = m.group(2).strip()
    
    # === IMPROVED LOCATION EXTRACTION (THIS IS THE KEY FIX) ===
    location_patterns = [
        r"(?i)(?:site location|project site)[\s:–\-–]*\s*([A-Za-z,\s\-–]+?)(?:,|\.|\n|Tamil Nadu|District|Vellore|India|$)",
        r"(?i)execution\s*at\s*[:\-]?\s*([A-Za-z\s\-–]+)",
        r"(?i)work\s*site\s*[:\-]?\s*([A-Za-z\s\-–]+)",
        r"(?i)delivery\s*at\s*[:\-]?\s*([A-Za-z\s\-–]+)",
        r"(?i)plant\s*location\s*[:\-]?\s*([A-Za-z\s\-–]+)",
    ]
    
    for pattern in location_patterns:
        m = re.search(pattern, text)
        if m:
            raw_location = m.group(1).strip()
            clean_location = re.sub(r'\s+', ' ', raw_location).title()
            info['location'] = clean_location
            break  # Stop at first match

    # Due Date
    m = re.search(r"(?i)(?:tender opening date|due date|submission date).*?([\d]{1,2}[\.\-\/][\d]{1,2}[\.\-\/][\d]{2,4})", text)
    if m: 
        info['due_date'] = m.group(1).strip()
    
    if info['due_date'] == 'N/A':
        m = re.search(r"(?i)due\s*date.*?[:\-]\s*([\d\/\-\.]+)", text)
        if m: 
            info['due_date'] = m.group(1).strip()
    
    return info

@st.cache_data(show_spinner=False)
def extract_content_for_topics(text):
    content = {title: [] for title in STANDARD_SLIDES}
    lines = [line.strip() for line in text.split('\n') if line.strip() and len(line) > 10]
    
    current_topic = None
    topic_keywords = {
        "Scope of Work": ["scope of work", "scope of tender", "complete manufacturing", "tender enquiry", "contract for"],
        "Pre-Qualification Criteria & Technical Capability": ["pre-qualification", "pqr", "technical capability"],
        "Technical Specifications": ["technical specification", "material specification", "tdc", "delivery conditions", "painting procedure"],
        "Commercial Terms": ["commercial terms", "price", "gst", "clauses"],
        "Documents Required": ["documents required", "documents to be submitted", "annexure"],
        "Eligibility Criteria": ["eligibility criteria", "pre-qualification", "pqr"],
        "Payment Terms": ["payment terms", "mode of payment", "eft", "rtgs"],
        "Security Deposit & Performance Bank Guarantee": ["security deposit", "performance bank guarantee", "pbg", "bank guarantee"],
        "Liquidated Damages": ["liquidated damages", "ld"],
        "Warranty & Guarantee": ["warranty", "guarantee", "defect liability"],
        "Safety Requirements": ["safety requirements", "safety appliances", "training"],
        "Bid Schedule & Key Dates": ["due date", "bid submission", "tender opening", "pre-bid", "validity"],
        "Delivery Period": ["delivery period", "completion period"],
        "Quality Standards": ["quality standards", "inspection", "quality control", "qap", "qwi"],
        "Evaluation Criteria": ["evaluation criteria"],
        "Expected Competition & Market Analysis": ["competition", "market"],
        "Bid Strategy & Recommendation": ["bid strategy"],
        "Critical Provisions & Compliance": ["critical provisions"],
        "Terms & Conditions Acceptance": ["terms & conditions"],
        "Make in India & Local Content": ["make in india", "local content", "class-i"],
        "Risk Analysis": ["risk analysis"],
        "Risk Mitigation Measures": ["risk mitigation"],
        "Financial Analysis": ["financial analysis"],
        "Organization & Team Structure": ["organization", "team structure"],
        "Quality Assurance Plan": ["quality assurance plan"],
        "Manufacturing Process Overview": ["manufacturing process"],
        "Inspection & Testing Protocol": ["inspection & testing"],
        "Packing & Transportation": ["packing", "transportation"],
        "Key Success Factors": ["key success factors"],
        "Conclusion & Recommendation": ["conclusion", "recommendation"],
    }
    
    buffer = []
    for line in lines:
        lower_line = line.lower()
        matched = False
        for topic, keywords in topic_keywords.items():
            if any(kw in lower_line for kw in keywords) and re.match(r'^[A-Z0-9\.\s\-:]+$', line) is None:
                if current_topic and buffer:
                    content[current_topic].append(' '.join(buffer))
                current_topic = topic
                buffer = [line]
                matched = True
                break
        if not matched and current_topic:
            buffer.append(line)
    
    if current_topic and buffer:
        content[current_topic].append(' '.join(buffer))
    
    for topic in content:
        content[topic] = [re.sub(r'\s+', ' ', item).strip() for item in content[topic] if len(item) > 20][:10]

    return content

# PPT functions
def add_wave_background(slide):
    wave = slide.shapes.add_shape(MSO_SHAPE.WAVE, Inches(-1), Inches(4.8), Inches(15), Inches(3))
    wave.fill.solid()
    wave.fill.fore_color.rgb = RGBColor(0, 176, 240)
    wave.line.fill.background()
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(0), Inches(3.5), Inches(7.5))
    tri.fill.solid()
    tri.fill.fore_color.rgb = RGBColor(173, 216, 230)
    tri.line.fill.background()

def add_logo(slide):
    box = slide.shapes.add_textbox(Inches(10.5), Inches(0.1), Inches(2.5), Inches(0.6))
    tf = box.text_frame
    tf.text = "TETAKISU"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 32, 96)
    p.alignment = PP_ALIGN.RIGHT

def add_content_slide(prs, title, bullets, use_table=False, subheadings=None):
    # Split bullets into chunks of max 10 per slide to avoid overflow
    max_bullets_per_slide = 10
    for page_idx, i in enumerate(range(0, len(bullets), max_bullets_per_slide)):
        part = bullets[i:i + max_bullets_per_slide]
        slide_title = title if page_idx == 0 else f"{title} (Contd. {page_idx + 1})"
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_wave_background(slide)
        add_logo(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.text = slide_title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

        top = Inches(1.6)

        if use_table:
            # Table layout (for key-value slides)
            rows = len(part) + 1
            table = slide.shapes.add_table(rows, 2, Inches(0.8), top, Inches(11.5), Inches(0.5 * rows))
            header = table.table.rows[0].cells
            header[0].text = "Key"
            header[1].text = "Value"
            for cell in header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 32, 96)
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.font.bold = True
            for j, (k, v) in enumerate(part, 1):
                table.table.cell(j, 0).text = k
                table.table.cell(j, 1).text = v if v != 'N/A' else '-'
        else:
            if subheadings and page_idx == 0:
                # Only show subheadings on first page
                for sub, subs in subheadings.items():
                    sub_box = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(0.5))
                    tf = sub_box.text_frame
                    tf.text = sub
                    tf.paragraphs[0].font.size = Pt(26)
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)
                    top += Inches(0.6)

                    # Split sub-bullets across pages if needed
                    sub_part = subs[i:i + max_bullets_per_slide] if sub == list(subheadings.keys())[1] else subs[:max_bullets_per_slide]
                    for b in sub_part:
                        if top > Inches(6.5):  # Prevent overflow
                            break
                        box = slide.shapes.add_textbox(Inches(1.2), top, Inches(11), Inches(0.45))
                        tf = box.text_frame
                        tf.text = f"• {b}"
                        tf.paragraphs[0].font.size = Pt(18)
                        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                        top += Inches(0.5)
            else:
                # Normal bullet list
                for b in part:
                    if top > Inches(6.8):  # Safety margin
                        break
                    box = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(0.45))
                    tf = box.text_frame
                    tf.text = f"• {b}"
                    # Auto-adjust font size if text is very long
                    if len(b) > 120:
                        tf.paragraphs[0].font.size = Pt(16)
                    elif len(b) > 90:
                        tf.paragraphs[0].font.size = Pt(17)
                    else:
                        tf.paragraphs[0].font.size = Pt(18)
                    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                    top += Inches(0.5)        
    
# NEW: Project Overview - Bullet style like your screenshot
def add_project_overview_bullet_slide(prs, tender_info, text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_background(slide)
    add_logo(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.text = "Project Overview"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

    # Auto-detect everything from the tender text
    security = "N/A"
    m = re.search(r"security deposit.*?(\d+%?)\s", text, re.I)
    if m: security = m.group(1).strip()

    pbg = "N/A"
    m = re.search(r"performance bank guarantee.*?(\d+%?)\s", text, re.I)
    if m: pbg = m.group(1).strip()

    validity = "N/A"
    m = re.search(r"bid validity.*?(\d+\s*days?)", text, re.I)
    if m: validity = m.group(1).strip()

    ld = "N/A"
    m = re.search(r"liquidated damages.*?(\d+\.?\d*\%?)\s", text, re.I)
    if m: ld = m.group(1).strip()

    warranty = "N/A"
    m = re.search(r"warranty.*?(\d+\s*(months?|years?))", text, re.I)
    if m: warranty = m.group(1).strip()

    completion = "N/A"
    m = re.search(r"completion period|delivery period.*?(\d+\s*days?)", text, re.I)
    if m: completion = m.group(1).strip()

    bullets = [
        f"• Client / Owner: {tender_info['client']}",
        f"• Site Location: {site_location}",
        f"• Scope of tender: {tender_info['tender_name']}",
        "• Funding: Own resources / commercial borrowings",
        f"• Tender Value in Rs.: N/A",
        f"• Security Deposit: {security} of Contract Value within 15 days of PO",
        f"• Performance Bank Guarantee: {pbg} performance Bank guarantee",
        f"• Bid Validity: {validity}",
        f"• LD: {ld} of contract price for each week of delay, up to a maximum of 10% of contract price",
        f"• Warranty: The supplied items must be guaranteed for a minimum of {warranty} against any defects in material, workmanship, or design",
        f"• Completion Period: {completion} of Purchase Order (PO) date and responsibilities at manufacturer's site after trial sub-assembly testing"
    ]

    top = Inches(1.8)
    for bullet in bullets:
        box = slide.shapes.add_textbox(Inches(0.8), top, Inches(12), Inches(0.5))
        tf = box.text_frame
        tf.text = bullet
        tf.paragraphs[0].font.size = Pt(20)
        top += Inches(0.55)



# NEW: Projects & their Locations - High Quality Static Map (Reliable & Beautiful)
def add_projects_locations_slide(prs, tender_info, coords, site_location):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_background(slide)
    add_logo(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = "Projects & their Locations"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

    # Add table
    rows = 2
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12)
    height = Inches(1)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(1.5)

    # Header
    headers = ["Project Name / Description", "Project Location", "Client / Owner", "Tender Reference"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        tf = cell.text_frame
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(12)

    # Data row (auto from tender)
    desc = tender_info.get('tender_name', 'N/A')
    loc = site_location if site_location != "Location Not Identified" else tender_info.get('location', 'N/A')
    client = tender_info.get('client', 'N/A')
    ref = tender_info.get('tender_no', 'N/A')
    data = [desc, loc, client, ref]
    for i, d in enumerate(data):
        cell = table.cell(1, i)
        cell.text = d
        tf = cell.text_frame
        tf.paragraphs[0].font.size = Pt(12)

    # Add static map image if coords available
    if coords:
        lat, lon = coords
        map_url = f"https://staticmap.openstreetmap.de/staticmap.php?center={lat},{lon}&zoom=12&size=800x400&maptype=mapnik&markers={lat},{lon},red-pushpin"
        try:
            resp = requests.get(map_url, timeout=10)
            if resp.status_code == 200:
                img_stream = BytesIO(resp.content)
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(3), width=Inches(12), height=Inches(3))
        except:
            pass  # Skip if fails

# NEW: Projects - Site Access
def add_site_access_slide(prs, infra):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_background(slide)
    add_logo(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.text = "Projects - Site Access"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

    table_data = []
    if infra.get("Airport"):
        table_data.append(["Nearest Airport", infra["Airport"][0]])
    if infra.get("Hotel"):
        table_data.append(["Nearest Hotel / Stay Option", infra["Hotel"][0]])
    if infra.get("Railway Station"):
        table_data.append(["Nearest Railway Station", infra["Railway Station"][0]])
    if infra.get("Bus Stand"):
        table_data.append(["Local Transport Options", infra["Bus Stand"][0]])
    
    if table_data:
        rows = len(table_data) + 1
        table = slide.shapes.add_table(rows, 2, Inches(0.8), Inches(1.8), Inches(11), Inches(0.5 * rows))
        hdr = table.table.rows[0].cells
        hdr[0].text = ""
        hdr[1].text = ""
        for cell in hdr:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 32, 96)
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.bold = True
        for i, row in enumerate(table_data, 1):
            table.table.cell(i, 0).text = row[0]
            table.table.cell(i, 1).text = row[1]
    else:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(0.5))
        tf = box.text_frame
        tf.text = "No infrastructure details detected"
        tf.paragraphs[0].font.size = Pt(20)

# NEW: Site Location Slide
def add_site_location_slide(prs, site_location, coords):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_background(slide)
    add_logo(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = title_box.text_frame
    tf.text = "Site Location & Details"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
    tf = content_box.text_frame
    tf.word_wrap = True
    bullets = [
        f"• Location: {site_location}",
        f"• Coordinates: {coords if coords else 'N/A'}",
        "• Address: N/A",
        "• Access: Well connected by road and rail",
        "• Infrastructure: Modern manufacturing facility with advanced equipment"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(24)

    # Add static map if coords
    if coords:
        lat, lon = coords
        map_url = f"https://staticmap.openstreetmap.de/staticmap.php?center={lat},{lon}&zoom=10&size=800x400&maptype=mapnik&markers={lat},{lon},red-pushpin"
        try:
            resp = requests.get(map_url, timeout=10)
            if resp.status_code == 200:
                img_stream = BytesIO(resp.content)
                slide.shapes.add_picture(img_stream, Inches(1), Inches(4.5), width=Inches(11), height=Inches(2.5))
        except:
            pass

def add_title_slide(prs, tender_info):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_wave_background(slide)
    add_logo(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11), Inches(1))
    tf = title_box.text_frame
    tf.text = "Tender Analysis Presentation"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(11), Inches(1))
    tf = subtitle_box.text_frame
    tf.text = f"Tender No: {tender_info['tender_no']}\nClient: {tender_info['client']}\nLocation: {tender_info['location']}"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 32, 96)

if uploaded_files:
    st.session_state.pop("geo_data", None)
    
    with st.spinner("Processing documents..."):
        all_text = ""
        all_tables = []  # To collect tables from ALL files

                # ------------------ IMPROVED FILE PROCESSING ------------------
        for uploaded_file in uploaded_files:
            file_bytes = uploaded_file.read()
            filename = uploaded_file.name.lower()
            
            # Quick magic number check (first 16 bytes) to reliably detect format
            header = file_bytes[:16]
            
            text = ""
            tables = []
            
            if header.startswith(b'%PDF-'):
                # It's a PDF (normal or scanned)
                text, tables = extract_text_and_tables_from_pdf(file_bytes, uploaded_file.name)
            
            elif header.startswith(b'\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1'):  # OLE compound → classic .xls, .doc
                if filename.endswith(('.xls', '.xlsx')):
                    text, tables = extract_text_and_tables_from_excel(file_bytes, uploaded_file.name)
                elif filename.endswith('.docx'):
                    text, tables = extract_text_and_tables_from_docx(file_bytes, uploaded_file.name)
                else:
                    st.warning(f"Could not process OLE file: {uploaded_file.name}")
            
            elif header.startswith(b'PK\x03\x04'):  # ZIP container → modern .xlsx, .docx
                if filename.endswith(('.xlsx', '.xls')):
                    text, tables = extract_text_and_tables_from_excel(file_bytes, uploaded_file.name)
                elif filename.endswith('.docx'):
                    text, tables = extract_text_and_tables_from_docx(file_bytes, uploaded_file.name)
                else:
                    st.warning(f"Could not process ZIP-based file: {uploaded_file.name}")
            
            else:
                st.warning(f"Unknown/unsupported file format: {uploaded_file.name}")
                text, tables = "", []
            
            all_text += f"\n\n--- {uploaded_file.name} ---\n{text}"
            all_tables.extend(tables)

        # Now that all files are processed → extract tender info
        tender_info = extract_tender_info(all_text)
        
        # Clean up N/A values dynamically (safe for any tender)
        for key in tender_info:
            value = tender_info.get(key, "N/A")
            if value == "N/A" or not value or value.strip() == "":
                if key.lower() in ["tender_value", "contract_value", "value"]:
                    tender_info[key] = "As per priced BOQ / final negotiation"
                elif key.lower() in ["ld", "liquidated_damages", "penalty"]:
                    tender_info[key] = "As per tender conditions (typically max 10%)"
                elif key.lower() in ["security_deposit", "performance_guarantee", "pbg", "bank_guarantee"]:
                    tender_info[key] = "As per GCC / tender terms"
                elif key.lower() in ["warranty", "guarantee"]:
                    tender_info[key] = "Minimum 12-18 months from supply or as specified"
                elif key.lower() in ["bid_validity", "validity"]:
                    tender_info[key] = "As per tender schedule"
                elif key.lower() in ["completion_period", "delivery", "period"]:
                    tender_info[key] = "As per tender schedule / PO terms"
                else:
                    tender_info[key] = "To be confirmed from tender / PO"
        
        content_by_topic = extract_content_for_topics(all_text)

        # Generate dynamic table summary (this is the director-friendly short version)
        table_summary = generate_table_summary(all_tables)

        # GEO DETECTION (your existing logic - unchanged)
        if "geo_data" not in st.session_state:
            site_location, coords = detect_location_and_coords(all_text, tender_info)
            climate = get_climate_data(*coords) if coords else {}
            infra = get_nearby_infra(*coords) if coords else {}
            st.session_state.geo_data = {
                "site_location": site_location, 
                "coords": coords,
                "climate": climate,
                "infra": infra,
                "table_summary": table_summary  # Optional: save for PPT later
            }
        else:
            geo = st.session_state.geo_data
            site_location, coords, climate, infra = geo["site_location"], geo["coords"], geo["climate"], geo["infra"]
            table_summary = geo.get("table_summary", "No tables found.")

        st.success(f"Processed {len(uploaded_files)} files | Tender: {tender_info['tender_no']}")

        # ---------------- PROFESSIONAL DYNAMIC SUMMARY DISPLAY ----------------
    st.markdown("## 📊 Summary of Key Quality & Inspection Tables")
    
    if all_tables:
        with st.expander("View detailed auto-generated summary from all tables", expanded=True):
            st.markdown(table_summary)
            st.caption(f"Dynamically extracted from {len(all_tables)} tables in {len(uploaded_files)} documents • Fully based on uploaded content")
    else:
        st.info("No tables found in the uploaded documents.")
    
    st.divider()

# ---------------- INTERACTIVE KEYWORD VIEW ----------------
if uploaded_files and content_by_topic:
    st.markdown("## 📌 Key Sections")
    left, right = st.columns([1, 3])
    with left:
        selected_topic = st.radio(
            "Select Keyword",
            list(content_by_topic.keys())
        )
    with right:
     st.markdown(f"## 🔍 {selected_topic}")
    
    # Get the raw extracted content for this topic
    raw_content = content_by_topic.get(selected_topic, [])
    
    # Join into one string and create short summary
    raw_text = " ".join(raw_content)
    clean_summary = summarize_text(raw_text, max_sentences=6)  # You can change 6 to 4 or 8
    
    # Show the nice short summary
    st.markdown(clean_summary)
    
    # Optional: Full raw text hidden in expander
    with st.expander("View Full Extracted Text (original)"):
        if raw_content:
            for item in raw_content:
                st.markdown(f"- {item}")
        else:
            st.markdown("No details found in this section.")
    
    st.divider()
    
    # Keep your existing site location, climate, infra display
    st.markdown("### 📍 Site Location")
    st.write(site_location)
    if coords:
        m = folium.Map(location=coords, zoom_start=10)
        folium.Marker(coords, popup=site_location).add_to(m)
        st_folium(m, height=250, key="map")
    st.divider()
    st.markdown("### 🌦 Climate Conditions")
    st.metric("Temperature", climate.get("temperature", "N/A"))
    st.metric("Humidity", climate.get("humidity", "N/A"))
    st.metric("Rainfall", climate.get("precipitation", "N/A"))
    st.divider()
    st.markdown("### 🏗 Nearby Infrastructure")
    for k, v in infra.items():
        st.markdown(f"**{k}**")
        if v:
            for item in v:
                st.markdown(f"- {item}")
        else:
            st.markdown("- No major facility available within defined radius")
else:
    st.info("Upload tender documents to view extracted sections.")

if st.button("Generate 60-70 Slide Professional PPT", type="primary"):
    with st.spinner("Creating director-ready presentation..."):
        table_summary = st.session_state.geo_data.get("table_summary", "No tables detected in the documents.")
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Updated Title Slide as per screenshot (fully automatic)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_wave_background(slide)
        add_logo(slide)

        # Client Name (auto)
        client_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(2))
        tf = client_box.text_frame
        tf.text = tender_info['client'].upper()
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Nature of Work (auto)
        nature_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(1))
        tf = nature_box.text_frame
        tf.text = f"Nature of Work: {tender_info['tender_name']}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Tender No. (auto)
        tender_no_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.8))
        tf = tender_no_box.text_frame
        tf.text = f"Tender No.: {tender_info['tender_no']}"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Bid / No-Bid Screening (standard phrase)
        bid_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(0.8))
        tf = bid_box.text_frame
        tf.text = "Bid / No-Bid Screening"
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Agenda Slide - Clean, no grey highlight
        agenda_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_wave_background(agenda_slide)
        add_logo(agenda_slide)

        agenda_title_box = agenda_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(11), Inches(1))
        tf = agenda_title_box.text_frame
        tf.text = "Agenda"
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

        top = Inches(1.8)
        items = [
            "Project Overview",
            "Key Tender Conditions",
            "Evaluation Criteria",
            "Expected Competition & Target Price",
            "Tender Organisation & Association of TETAKISU Representatives",
            "Bid Strategy & Recommendation"
        ]

        for item in items:
            box = agenda_slide.shapes.add_textbox(Inches(1.0), top, Inches(11), Inches(0.5))
            tf = box.text_frame
            tf.text = item
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = False
            top += Inches(0.8)

        # NEW SLIDES AFTER AGENDA
        add_project_overview_bullet_slide(prs, tender_info, text)
        add_projects_locations_slide(prs, tender_info, coords, site_location)
        add_site_access_slide(prs, infra)

        # Regular content slides
        for title in STANDARD_SLIDES[1:]:
            if title in ["Project Overview", "Site Location & Details", "Climatic Conditions Summary", "Nearby Infrastructure"]:
                continue
            bullets = content_by_topic.get(title, ["No specific details found."])
            subheadings = None
            if title in ["Scope of Work", "Payment Terms", "Bid Schedule & Key Dates"]:
                subheadings = {"Key Details": bullets[:6], "Additional Notes": bullets[6:]}  # Better split
            add_content_slide(prs, title, bullets, subheadings=subheadings)

        # Rich slides (you can keep or remove duplicate Site Location)
        add_site_location_slide(prs, site_location, coords)

        # ---- NEW SLIDE: Quality Tables Summary ----
        summary_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        add_wave_background(summary_slide)
        add_logo(summary_slide)

                       # ---- ENHANCED ILLUSTRATIONS SLIDE WITH REAL EMBEDDED IMAGES ----
        illus_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_wave_background(illus_slide)
        add_logo(illus_slide)

        title_box = illus_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        tf.text = "Key Component & Process Illustrations"
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

        tender_lower = all_text.lower()
        image_data = []  # (url, caption)

        if any(word in tender_lower for word in ["esp", "electrostatic precipitator", "precipitator"]):
            image_data = [
                ("https://upload.wikimedia.org/wikipedia/commons/thumb/4/4f/Electrostatic_precipitator_principle.svg/800px-Electrostatic_precipitator_principle.svg.png", "ESP Working Principle & Configuration"),
                ("https://upload.wikimedia.org/wikipedia/commons/thumb/8/8f/ESP_Collecting_plates.jpg/800px-ESP_Collecting_plates.jpg", "Collecting Plates"),
                ("https://www.babcock.com/home/Images/ESP_Rapping_System.jpg", "Rapping System for Ash Removal")
            ]
        elif any(word in tender_lower for word in ["boiler", "pressure parts", "drum", "header"]):
            image_data = [
                ("https://upload.wikimedia.org/wikipedia/commons/thumb/0/0e/Boiler_pressure_parts.jpg/800px-Boiler_pressure_parts.jpg", "Boiler Pressure Parts Layout"),
                ("https://www.babcock.com/home/Images/Steam_Drum.jpg", "Steam Drum Cross Section")
            ]
        elif any(word in tender_lower for word in ["pfbr", "fast breeder", "kalpakkam", "reactor"]):
            image_data = [
                ("https://upload.wikimedia.org/wikipedia/commons/thumb/9/9f/Fast_Breeder_Reactor_Schematic.svg/800px-Fast_Breeder_Reactor_Schematic.svg.png", "PFBR Reactor Vessel & Core Assembly")
            ]
        elif any(word in tender_lower for word in ["painting", "surface preparation", "blast", "dft"]):
            image_data = [
                ("https://upload.wikimedia.org/wikipedia/commons/thumb/5/5e/Sa_2.5_blast_cleaning.jpg/800px-Sa_2.5_blast_cleaning.jpg", "Sa 2½ Blast Cleaning Process"),
                ("https://www.elcometer.com/images/stories/products/coating-thickness-gauges/elcometer-456.jpg", "DFT Measurement with Elcometer")
            ]
        else:
            image_data = [
                ("https://upload.wikimedia.org/wikipedia/commons/thumb/7/7f/Heavy_fabrication_process.jpg/800px-Heavy_fabrication_process.jpg", "Heavy Structural Fabrication & Assembly")
            ]

        # Bullets
        content_box = illus_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.text = "Relevant illustrations based on tender scope:"
        for _, caption in image_data:
            p = tf.add_paragraph()
            p.text = f"• {caption}"
            p.font.size = Pt(24)

        # Embed images (2 per row, small size)
        left = Inches(1)
        right = Inches(7)
        top = Inches(5)
        img_w = Inches(5)
        img_h = Inches(2.5)

        for i, (url, _) in enumerate(image_data[:4]):  # Max 4 images
            try:
                resp = requests.get(url, timeout=10)
                if resp.status_code == 200:
                    stream = BytesIO(resp.content)
                    x = left if i % 2 == 0 else right
                    y = top if i < 2 else top + img_h + Inches(0.3)
                    illus_slide.shapes.add_picture(stream, x, y, width=img_w, height=img_h)
            except:
                pass  # Skip if image fails to load

        # Note
        note_box = illus_slide.shapes.add_textbox(Inches(1), Inches(7), Inches(11), Inches(0.6))
        tf = note_box.text_frame
        tf.text = "Note: Actual components to be manufactured as per approved drawings"
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

        # Title
        title_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        tf = title_box.text_frame
        tf.text = "Quality Tables Summary"
        tf.paragraphs[0].font.size = Pt(44)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0, 176, 240)

        # Content (bullet points from dynamic summary)
        content_box = summary_slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.text = table_summary
        tf.paragraphs[0].level = 0  # In case first line isn't bulleted
        for paragraph in tf.paragraphs:
            if paragraph.text.startswith("•"):
                paragraph.level = 0
            paragraph.font.size = Pt(24)
            paragraph.font.name = "Calibri"
            paragraph.space_after = Pt(12)  # Ensure proper line breaks
        for para in tf.paragraphs:
            para.font.size = Pt(24)
            para.font.name = "Calibri"
            para.space_after = Pt(10)
            
        bio = io.BytesIO()
        prs.save(bio)
        bio.seek(0)

        st.download_button(
            "⬇️ Download Professional PPT",
            bio,
            f"Tender_{tender_info['tender_no']}_{datetime.now().strftime('%Y%m%d')}.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        st.success(f"Generated professional PPT with your exact layout and embedded map!")

else:
    st.info("Upload tender documents to generate a professional presentation.")
